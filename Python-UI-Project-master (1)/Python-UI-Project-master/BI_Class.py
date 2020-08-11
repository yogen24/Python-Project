import numpy as np
import pandas as pd
from datetime import datetime

#Importing PPTX, defining layout and adding slide.
from pptx import Presentation

#Importing libraries to plot the chart
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt

#Reading, Cleaning District.csv file
district=pd.read_csv('district.csv')
district=district.drop(['A5','A6','A7','A8','A9','A10','A12','A13','A14','A15','A16'],axis=1)
district=district.rename(columns={'A1':'district_id','A2':'District_Name','A3':'Region','A4':'Population','A11':'Avg_Salary'})

#Reading, Cleaning Account.csv file
account=pd.read_csv('account1.csv')
account['frequency'].replace({"POPLATEK MESICNE": "Monthly Issuance", "POPLATEK TYDNE": "Weekly Issuance",
                             'POPLATEK PO OBRATU':'Issuance After Transaction'}, inplace=True)
account['Account_date'] = pd.to_datetime(account['Modified_date'].astype(str), format='%Y%m%d')
account=account.drop(['date','Modified_date'],axis=1)

#Reading, Cleaning Order.csv file
order=pd.read_csv('order1.csv')
order['k_symbol'].replace({"POJISTNE": "Insurance Payment", "SIPO": "Household Payment",'LEASING':'Leasing',
                             'UVER':'Loan Payment',' ':'Other Payments'}, inplace=True)
order=order.rename(columns={'k_symbol':'Order_Description'})
order=order.drop(['account_to'],axis=1)

#Reading, Cleaning Trans.csv file
trans=pd.read_csv('trans1.csv')
trans['type'].replace({"PRIJEM": "Credit", "VYDAJ": "Withdrawal"},inplace=True)
trans['operation'].replace({"VYBER KARTOU": "Credit Card Withdrawal", "VKLAD": "Credit in Cash",
                            'PREVOD Z UCTU':'Collection from Another Bank', "VYBER":"Withdrawal in Cash",
                             'PREVOD NA UCET':'Remittance to Another Bank'}, inplace=True)
trans['operation'].fillna('Other transactions',inplace=True)
trans=trans.drop(['delete','balance','account','k_symbol'],axis=1)
trans=trans.rename(columns={'date':'Trans_date','type':'trans_type','amount':'trans_amount'})
trans['Trans_date']=pd.to_datetime(trans['Trans_date'])
trans.set_index (trans['Trans_date'], inplace = True)
trans=trans.sort_index()

#Adding yearly inflow and outflow columns to TRANS table
trans= trans.assign (positive = np.nan)
trans= trans.assign (negative= np.nan)

trans['positive'] = trans['trans_amount'].where (trans['operation'] == 'Credit in Cash', trans['positive'] )
trans['positive'] = trans['trans_amount'].where (trans['operation'] == 'Collection from Another Bank', trans['positive'])

trans['negative'] = trans['trans_amount'].where (trans['operation'] == 'Withdrawal in Cash', trans['negative'] )
trans['negative'] = trans['trans_amount'].where (trans['operation'] == 'Remittance to Another Bank', trans['negative'] )
trans['negative'] = trans['trans_amount'].where (trans['operation'] == 'Credit Card Withdrawal', trans['negative'])

#Separating date on basis of year and month_year
trans['Year']=trans['Trans_date'].apply(lambda i:i.year)
trans['month_year'] = trans['Trans_date'].apply(lambda x: x.strftime('%b-%y')) 

def trans_yearly():
    prs=Presentation()
    layout=prs.slide_layouts[6]
    from pptx.util import Pt
    #Description slide
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)

    left = width = height = Inches(1)
    top= Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "                 Description\n"
    p.font.bold = True
    p.font.size = Pt(40)

    p = tf.add_paragraph()
    p.text = "            Dimension: Bank      Granularity: Years\n"
    p.font.size = Pt(25)

    p = tf.add_paragraph()
    p.text = "                        Data used:  Transactions.csv\n\n"
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "The following chart shows the variation of Total Fund Flow (Inflow & Outflow)\nfor all banks over each year from 1993 to 1998. \n\nAnd the following table shows the corresponding fund totals for each bank."
    p.font.size = Pt(20)

    inflowsY = trans['positive'].resample('Y').sum ()
    outflowsY = trans['negative'].resample('Y').sum ()
    
    slide=prs.slides.add_slide(layout)
    mydata=CategoryChartData()

    #Inputing Years as dimension
    mydata.categories=np.sort(trans['Year'].unique())

    #Inputing Total inflow and outflow amount as metric
    mydata.add_series('Fund Inflow',inflowsY.values/1000000)
    mydata.add_series('Fund Outflow',outflowsY.values/1000000)
    
    x,y,cx,cy=Inches(1),Inches(0.5),Inches(9),Inches(5.2)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.LINE,x,y,cx,cy,mydata).chart

    chart.has_legend=True
    from pptx.util import Pt
    chart.chart_title.has_text_frame=True
    chart.chart_title.text_frame.text='Banks Yearly Fund Flow (Millions USD)'
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    value_axis_title = chart.value_axis.axis_title
    value_axis_title.text_frame.text = "Amount (Millions USD)"

    x,y,cx,cy=Inches(0),Inches(6),Inches(10),Inches(1.2)
    shape=slide.shapes.add_table(3,14,x,y,cx,cy)
    table=shape.table

    cell1=table.cell(0,0)
    cell1.text='Bank'
    cell1.text_frame.paragraphs[0].font.size=Pt(14)
    cell2=table.cell(1,0)
    cell2.text='Total Inflow'
    cell2.text_frame.paragraphs[0].font.size=Pt(14)
    cell3=table.cell(2,0)
    cell3.text='Total Outflow'
    cell3.text_frame.paragraphs[0].font.size=Pt(13)
    from pptx.dml.color import RGBColor
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = RGBColor(185, 70, 70)
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = RGBColor(235,224,224)
    cell3.fill.solid()
    cell3.fill.fore_color.rgb = RGBColor(235,224,224)
    table.columns[0].width = Inches(1.4)
    cnt=1

    for i in trans.groupby('bank')['positive'].sum().index:
        cell = table.cell(0, cnt)
        cell.text = i
        cell2 = table.cell(1, cnt)
        cell2.text = str(round((trans.groupby('bank')['positive'].sum().values[cnt-1]/1000000))).split('.')[0]+'M'
        cell3 = table.cell(2, cnt)
        cell3.text = str(round((trans.groupby('bank')['negative'].sum().values[cnt-1]/1000000))).split('.')[0]+'M'    
        cnt=cnt+1

    prs.save('Yearly_Transactions.pptx')

def trans_monthly():
    prs=Presentation()
    layout=prs.slide_layouts[6]
    from pptx.util import Pt
    #Description slide
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)

    left = width = height = Inches(1)
    top= Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "                 Description\n"
    p.font.bold = True
    p.font.size = Pt(40)

    p = tf.add_paragraph()
    p.text = "            Dimension: Bank      Granularity: Months\n"
    p.font.size = Pt(25)

    p = tf.add_paragraph()
    p.text = "                        Data used:  Transactions.csv\n\n"
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "The following chart shows the variation of Total Fund Flow (Inflow & Outflow)\nfor all banks over each month from Jan-93 to Dec-98. \n\nAnd the following table shows the corresponding fund averages for each\nmonth across different years."
    p.font.size = Pt(20)

    #Adding monthly inflow and outflow columns to TRANS table
    inflowsM = trans['positive'].resample('M').sum ()
    outflowsM = trans['negative'].resample('M').sum ()

    slide=prs.slides.add_slide(layout)
    mydata=CategoryChartData()

    #Inputing Years as dimension
    mydata.categories=sorted(trans['month_year'].unique().tolist(),key=lambda date:datetime.strptime(date,'%b-%y'))

    #Inputing Total inflow and outflow amount as metric
    mydata.add_series('Fund Inflow',inflowsM.values/1000000)
    mydata.add_series('Fund Outflow',outflowsM.values/1000000)
    
    x,y,cx,cy=Inches(1),Inches(0.5),Inches(9),Inches(5.2)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.LINE,x,y,cx,cy,mydata).chart

    chart.has_legend=True
    from pptx.util import Pt
    chart.chart_title.has_text_frame=True
    chart.chart_title.text_frame.text='Banks Monthly Fund Flow (Millions USD)'
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    value_axis_title = chart.value_axis.axis_title
    value_axis_title.text_frame.text = "Amount (Millions USD)"
    
    x,y,cx,cy=Inches(0),Inches(6),Inches(10),Inches(1.2)
    shape=slide.shapes.add_table(3,13,x,y,cx,cy)
    table=shape.table

    cell1=table.cell(0,0)
    cell1.text='Month'
    cell1.text_frame.paragraphs[0].font.size=Pt(14)
    cell2=table.cell(1,0)
    cell2.text='Avg Inflow'
    cell2.text_frame.paragraphs[0].font.size=Pt(14)
    cell3=table.cell(2,0)
    cell3.text='Avg Outflow'
    cell3.text_frame.paragraphs[0].font.size=Pt(14)
    from pptx.dml.color import RGBColor
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = RGBColor(185, 70, 70)
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = RGBColor(235,224,224)
    cell3.fill.solid()
    cell3.fill.fore_color.rgb = RGBColor(235,224,224)
    table.columns[0].width = Inches(1.4)
    trans['month'] = pd.DatetimeIndex(trans['Trans_date']).month
    import calendar
    trans['month_name'] = trans['month'].apply(lambda x: calendar.month_name[x][0:3])

    x=trans.groupby('month_year')[['positive','negative']].resample('M').sum().reset_index()
    x['just_month']=x['month_year'].apply(lambda i:i.split('-')[0])
    y=x.groupby('just_month')['positive'].mean()
    z=x.groupby('just_month')['negative'].mean()
    fin=sorted(y.index.tolist(),key=lambda date:datetime.strptime(date,'%b'))

    cnt=1
    for i in fin:
        cell = table.cell(0, cnt)
        cell.text = i
        cell2 = table.cell(1, cnt)
        cell2.text = str(round((y[fin].values[cnt-1]/1000000))).split('.')[0]+'M'
        cell3 = table.cell(2, cnt)
        cell3.text = str(round((z[fin].values[cnt-1]/1000000))).split('.')[0]+'M'
        cnt=cnt+1

    prs.save('Monthly_Transactions.pptx')
    
def region_orders():
    prs=Presentation()
    layout=prs.slide_layouts[6]
    from pptx.util import Pt
    #Description slide
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)

    left = width = height = Inches(1)
    top= Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "                 Description\n"
    p.font.bold = True
    p.font.size = Pt(40)

    p = tf.add_paragraph()
    p.text = "            Dimension: Region      Granularity: YTD\n"
    p.font.size = Pt(25)

    p = tf.add_paragraph()
    p.text = "                Data used:  Districts.csv  &  Orders.csv\n\n"
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "The following chart shows the distribution of payments made by clients\n (type of payment and payment amount) in all the regions.\n\nAnd the corresponding table shows the sum of the payment types across\nthese regions"
    p.font.size = Pt(20)


    #Joining the order and district table
    df1=pd.merge(order,account,how='inner',on='account_id')
    df=pd.merge(df1,district,how='inner',on='district_id')
    df['Region']=df['Region'].apply(lambda i:i.title())

    #Grouping by region and description in ORDER table
    grp=df.groupby(['Region','Order_Description'])['amount'].sum()

    layout=prs.slide_layouts[6]
    slide=prs.slides.add_slide(layout)
    mydata=CategoryChartData()

    #Inputing various regions as dimension
    mydata.categories=np.sort(df['Region'].unique())

    #Inputing order description amounts as metric via loop
    for i in df['Order_Description'].unique():
        mydata.add_series(i,grp[grp.index.get_level_values('Order_Description') == i].values)
    
    x,y,cx,cy=Inches(0.5),Inches(0.5),Inches(9.5),Inches(5.2)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED,x,y,cx,cy,mydata).chart

    chart.has_legend=True
    from pptx.enum.chart import XL_LEGEND_POSITION
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    from pptx.util import Pt
    category_axis_title = chart.category_axis.axis_title
    category_axis_title.text_frame.text = "Region"
    value_axis_title = chart.value_axis.axis_title
    value_axis_title.text_frame.text = "Amount (USD)"
    chart.chart_title.has_text_frame=True
    chart.chart_title.text_frame.text='Total Amount(USD) by Payment Type for Each Region'
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    x,y,cx,cy=Inches(0.5),Inches(6),Inches(9),Inches(1)
    shape=slide.shapes.add_table(2,6,x,y,cx,cy)
    table=shape.table

    cell1=table.cell(0,0)
    cell1.text='Payment Type'
    cell2=table.cell(1,0)
    cell2.text='Total Amount'
    from pptx.dml.color import RGBColor
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = RGBColor(185, 70, 70)
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = RGBColor(235,224,224)

    cnt=1
    for i in df.groupby('Order_Description')['amount'].sum().index:
        cell = table.cell(0, cnt)
        cell.text = i
        cell2 = table.cell(1, cnt)
        cell2.text = str(round((df.groupby('Order_Description')['amount'].sum().values[cnt-1]/1000))).split('.')[0]+' K'
        cnt=cnt+1

    prs.save('Region_Orders.pptx')

def district_orders():
    prs=Presentation()
    layout=prs.slide_layouts[6]
    from pptx.util import Pt
    #Description slide
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)

    left = width = height = Inches(1)
    top= Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "                 Description\n"
    p.font.bold = True
    p.font.size = Pt(40)

    p = tf.add_paragraph()
    p.text = "            Dimension: District      Granularity: YTD\n"
    p.font.size = Pt(25)

    p = tf.add_paragraph()
    p.text = "                Data used:  Districts.csv  &  Orders.csv\n\n"
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "The following chart shows the distribution of payments made by clients\n (type of payment and payment amount) in the 10 most active districts.\n\nAnd the corresponding table shows the sum of the payment types across\nthese districs"
    p.font.size = Pt(20)

    #Joining the order and district table
    df1=pd.merge(order,account,how='inner',on='account_id')
    df=pd.merge(df1,district,how='inner',on='district_id')
    df['Region']=df['Region'].apply(lambda i:i.title())

    layout=prs.slide_layouts[6]
    slide=prs.slides.add_slide(layout)
    mydata=CategoryChartData()

    top=np.sort(df.groupby('District_Name')['amount'].sum().sort_values(ascending=False).iloc[0:10].index)
    topv=df[df['District_Name'].apply(lambda i:i in top)]
    grp2=topv.groupby(['District_Name','Order_Description'])['amount'].sum()

    #Inputing various districts as dimension
    mydata.categories=top

    #Inputing order description amounts as metric via loop
    for i in df['Order_Description'].unique():
        mydata.add_series(i,grp2[grp2.index.get_level_values('Order_Description') == i].values)
    
    x,y,cx,cy=Inches(0.5),Inches(0.5),Inches(9.5),Inches(5.2)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED,x,y,cx,cy,mydata).chart

    chart.has_legend=True
    from pptx.enum.chart import XL_LEGEND_POSITION
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    from pptx.util import Pt
    category_axis_title = chart.category_axis.axis_title
    category_axis_title.text_frame.text = "District"
    value_axis_title = chart.value_axis.axis_title
    value_axis_title.text_frame.text = "Amount (USD)"
    chart.chart_title.has_text_frame=True
    chart.chart_title.text_frame.text='Total Amount(USD) by Payment Type for Top 10 Districts'
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    x,y,cx,cy=Inches(0.5),Inches(6),Inches(9),Inches(1)
    shape=slide.shapes.add_table(2,6,x,y,cx,cy)
    table=shape.table

    cell1=table.cell(0,0)
    cell1.text='Payment Type'
    cell2=table.cell(1,0)
    cell2.text='Total Amount'
    from pptx.dml.color import RGBColor
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = RGBColor(185, 70, 70)
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = RGBColor(235,224,224)

    cnt=1
    for i in df.groupby('Order_Description')['amount'].sum().index:
        cell = table.cell(0, cnt)
        cell.text = i
        cell2 = table.cell(1, cnt)
        cell2.text =str(round((grp2[grp2.index.get_level_values('Order_Description') == i].values.sum()/1000))).split('.')[0]+' K'
        cnt=cnt+1

    prs.save('District_Orders.pptx')
    
def bank_orders():
    prs=Presentation()
    layout=prs.slide_layouts[6]
    from pptx.util import Pt
    #Description slide
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)

    left = width = height = Inches(1)
    top= Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "                 Description\n"
    p.font.bold = True
    p.font.size = Pt(40)

    p = tf.add_paragraph()
    p.text = "            Dimension: Bank      Granularity: YTD\n"
    p.font.size = Pt(25)

    p = tf.add_paragraph()
    p.text = "                           Data used:  Orders.csv\n\n"
    p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "The following chart shows the distribution of payments made by clients\n (type of payment and payment amount) to all the banks.\n\nAnd the corresponding table shows the sum of the payment types across\nthese banks"
    p.font.size = Pt(20)


    #Grouping by bank and description in ORDER table
    grp=order.groupby(['bank_to','Order_Description'])['amount'].sum()

    layout=prs.slide_layouts[6]
    slide=prs.slides.add_slide(layout)
    mydata=CategoryChartData()

    #Inputing various banks as dimension
    mydata.categories=np.sort(order['bank_to'].unique())

    #Inputing order description amounts as metric via loop
    for i in order['Order_Description'].unique():
        mydata.add_series(i,grp[grp.index.get_level_values('Order_Description') == i].values)
    
    x,y,cx,cy=Inches(1),Inches(0.3),Inches(9),Inches(5.7)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED,x,y,cx,cy,mydata).chart

    chart.has_legend=True
    from pptx.enum.chart import XL_LEGEND_POSITION
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    from pptx.util import Pt
    category_axis_title = chart.category_axis.axis_title
    category_axis_title.text_frame.text = "Bank Name"
    value_axis_title = chart.value_axis.axis_title
    value_axis_title.text_frame.text = "Amount (USD)"
    chart.chart_title.has_text_frame=True
    chart.chart_title.text_frame.text='Total Amount(USD) by Payment Type for Each Bank'
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    x,y,cx,cy=Inches(0.5),Inches(6),Inches(9),Inches(1)
    shape=slide.shapes.add_table(2,6,x,y,cx,cy)
    table=shape.table

    cell1=table.cell(0,0)
    cell1.text='Payment Type'
    cell2=table.cell(1,0)
    cell2.text='Total Amount'
    from pptx.dml.color import RGBColor
    cell1.fill.solid()
    cell1.fill.fore_color.rgb = RGBColor(185, 70, 70)
    cell2.fill.solid()
    cell2.fill.fore_color.rgb = RGBColor(235,224,224)

    cnt=1
    for i in order.groupby('Order_Description')['amount'].sum().index:
        cell = table.cell(0, cnt)
        cell.text = i
        cell2 = table.cell(1, cnt)
        cell2.text = str(round(order.groupby('Order_Description')['amount'].sum().values[cnt-1]/1000)).split('.')[0]+' K'
        cnt=cnt+1

    prs.save('Bank_Orders.pptx')